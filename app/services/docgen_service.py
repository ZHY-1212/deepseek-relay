"""Document generation using AI + python-pptx/openpyxl/python-docx."""

import io
import httpx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxColor

from app.config import settings


async def generate_doc(doc_type: str, prompt: str) -> tuple[bytes, str]:
    """Generate a document. Returns (file_bytes, filename)."""

    # Step 1: Ask AI to generate structured content
    system_prompts = {
        "pptx": "你是一个PPT生成助手。根据用户需求，生成一个JSON格式的PPT大纲。格式：{\"title\":\"总标题\",\"slides\":[{\"title\":\"页标题\",\"points\":[\"要点1\",\"要点2\",\"要点3\"]}]}。输出纯JSON，不要markdown代码块。",
        "xlsx": "你是一个Excel生成助手。根据用户需求，生成一个JSON格式的表格数据。格式：{\"title\":\"表名\",\"headers\":[\"列1\",\"列2\"],\"rows\":[[\"数据1\",\"数据2\"],[\"数据3\",\"数据4\"]]}。输出纯JSON，不要markdown代码块。",
        "docx": "你是一个Word文档生成助手。根据用户需求，生成一个JSON格式的文档大纲。格式：{\"title\":\"标题\",\"sections\":[{\"heading\":\"小标题\",\"body\":\"段落内容...\"}]}。输出纯JSON，不要markdown代码块。",
    }

    system_prompt = system_prompts.get(doc_type, system_prompts["pptx"])

    async with httpx.AsyncClient(timeout=60) as client:
        resp = await client.post(
            "https://api.deepseek.com/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {settings.deepseek_api_key}",
                "Content-Type": "application/json",
            },
            json={
                "model": "deepseek-chat",
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt},
                ],
                "temperature": 0.7,
                "max_tokens": 4096,
            },
        )
        data = resp.json()
        raw = data["choices"][0]["message"]["content"].strip()
        # Remove markdown code block if present
        if raw.startswith("```"):
            raw = raw.split("\n", 1)[1]
            if raw.endswith("```"):
                raw = raw[:-3]
        import json
        try:
            content = json.loads(raw)
        except json.JSONDecodeError:
            content = {"title": prompt, "slides": [{"title": "AI生成", "points": [raw[:500] if raw else "请重试", "内容过长已截断"]}]}
            if doc_type == "xlsx":
                content = {"title": prompt, "headers": ["内容"], "rows": [[raw[:500] if raw else "请重试"]]}
            elif doc_type == "docx":
                content = {"title": prompt, "sections": [{"heading": "AI生成", "body": raw[:1000] if raw else "请重试"}]}

    # Step 2: Convert to file
    if doc_type == "pptx":
        return _make_pptx(content), f"{content.get('title','document')}.pptx"
    elif doc_type == "xlsx":
        return _make_xlsx(content), f"{content.get('title','document')}.xlsx"
    elif doc_type == "docx":
        return _make_docx(content), f"{content.get('title','document')}.docx"
    else:
        raise ValueError(f"Unknown doc type: {doc_type}")


def _make_pptx(data):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_pptx_bg(slide, RGBColor(0x63, 0x66, 0xF1))
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "Untitled")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = 1  # center

    # Content slides
    colors = [RGBColor(0x63, 0x66, 0xF1), RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0x06, 0xB6, 0xD4), RGBColor(0xF5, 0x9E, 0x0B)]
    for i, s in enumerate(data.get("slides", [])):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_pptx_bg(slide, colors[i % len(colors)])
        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(1.2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = s.get("title", "")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Points
        txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.5), Inches(4.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for j, pt in enumerate(s.get("points", [])):
            if j == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = f"  {pt}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.space_after = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_pptx_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _make_xlsx(data):
    wb = Workbook()
    ws = wb.active
    ws.title = data.get("title", "Sheet1")[:31]

    # Title row
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(data.get("headers", [])) or 1)
    title_cell = ws.cell(row=1, column=1, value=data.get("title", ""))
    title_cell.font = Font(size=18, bold=True, color="FFFFFF")
    title_cell.fill = PatternFill(start_color="6366F1", end_color="6366F1", fill_type="solid")
    title_cell.alignment = Alignment(horizontal="center")
    ws.row_dimensions[1].height = 40

    # Headers
    headers = data.get("headers", [])
    for i, h in enumerate(headers):
        cell = ws.cell(row=3, column=i + 1, value=h)
        cell.font = Font(size=12, bold=True, color="FFFFFF")
        cell.fill = PatternFill(start_color="8B5CF6", end_color="8B5CF6", fill_type="solid")
        cell.alignment = Alignment(horizontal="center")
    ws.row_dimensions[3].height = 30

    # Data rows
    for r, row in enumerate(data.get("rows", [])):
        for c, val in enumerate(row):
            cell = ws.cell(row=r + 4, column=c + 1, value=str(val))
            cell.font = Font(size=11)
            cell.alignment = Alignment(horizontal="left")
            if r % 2 == 0:
                cell.fill = PatternFill(start_color="F3F4F6", end_color="F3F4F6", fill_type="solid")

    # Auto-width
    for col in ws.columns:
        max_len = 0
        for cell in col:
            if cell.value:
                max_len = max(max_len, len(str(cell.value)))
        ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_docx(data):
    doc = Document()
    # Title
    title = doc.add_heading(data.get("title", "Untitled"), level=0)
    for run in title.runs:
        run.font.color.rgb = DocxColor(0x63, 0x66, 0xF1)

    for section in data.get("sections", []):
        heading = section.get("heading", "")
        body = section.get("body", "")
        if heading:
            h = doc.add_heading(heading, level=2)
            for run in h.runs:
                run.font.color.rgb = DocxColor(0x8B, 0x5C, 0xF6)
        if body:
            p = doc.add_paragraph(body)
            for run in p.runs:
                run.font.size = DocxPt(12)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()
